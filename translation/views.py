from rest_framework import generics
from .models import UploadedFile
from .serializers import UploadedFileSerializer
from django.views.generic import TemplateView
from django.http import FileResponse
from django.shortcuts import get_object_or_404
from django.views import View
from .models import UploadedFile


class UploadedFileListCreateView(generics.ListCreateAPIView):
    queryset = UploadedFile.objects.all()
    serializer_class = UploadedFileSerializer


class FileListView(TemplateView):
    template_name = 'file_list.html'

    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['files'] = UploadedFile.objects.all()
        return context


class FileDownloadView(View):
    def get(self, request, file_id):
        uploaded_file = get_object_or_404(UploadedFile, id=file_id)
        response = FileResponse(uploaded_file.file)
        response['Content-Disposition'] = f'attachment; filename="{uploaded_file.file.name}"'
        return response


download_file = FileDownloadView.as_view()



# translator/views.py
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework.permissions import IsAuthenticated
from .serializers import TranslatedFileSerializer
from .models import TranslatedFile
from googletrans import Translator
from pptx import Presentation
from openpyxl import load_workbook
from openpyxl.styles import Font
from openpyxl.utils import get_column_letter
from django.urls import reverse
import openpyxl
from googletrans import Translator
import os
from django.conf import settings
from pptx.util import Pt
from django.http import Http404




class TranslateFileView(APIView):
    queryset = TranslatedFile.objects.all()
    serializer_class = TranslatedFileSerializer

    def post(self, request):
        serializer = TranslatedFileSerializer(data=request.data)
        if serializer.is_valid():
            original_file = serializer.validated_data['original_file']
            target_language = serializer.validated_data['target_language']

            file_extension = original_file.name.split('.')[-1].lower()
            if file_extension == 'xlsx':
                translated_file_path = translate_xlsx(original_file, target_language)
            elif file_extension == 'pptx':
                translated_file_path = translate_pptx(original_file, target_language)
            else:
                return Response({'error': 'Unsupported file type'}, status=400)

            translated_file_instance = TranslatedFile.objects.create(
                original_file=original_file,
                # translated_file=translated_file_path,
                target_language=target_language,
            )
             # Provide a link to download the translated file
            download_link = f'/download/{translated_file_instance.id}/'
            return Response({'id': translated_file_instance.id, 'download_link': download_link})
        else:
            return Response(serializer.errors, status=400)


def translate_and_preserve_formatting(prs, target_language):
    formatted_text_data = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.text
                font = shape.text.font.name
                size = shape.text.font.size
                color = shape.text.font.color.rgb
                formatted_text_data.append({
                    'text': text,
                    'font': font,
                    'size': size,
                    'color': color
                })

    translator = Translator()
    translated_text_data = []
    for item in formatted_text_data:
        translated_text = translator.translate(item['text'], dest=target_language).text
        translated_text_data.append({
            'text': translated_text,
            'font': item['font'],
            'size': item['size'],
            'color': item['color']
        })

    i = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text.text = translated_text_data[i]['text']
                shape.text.font.name = translated_text_data[i]['font']
                shape.text.font.size = Pt(translated_text_data[i]['size'])
                shape.text.font.color.rgb = translated_text_data[i]['color']
                i += 1



def translate_xlsx(file_obj, target_language):
    try:
        workbook = load_workbook(file_obj, read_only=True)
    except Exception as e:
        return str(e)

    translator = Translator()

    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value is not None and isinstance(cell.value, str):
                    translated_text = translator.translate(cell.value, dest=target_language).text
                    cell.value = translated_text

    translated_file_path = os.path.join(settings.MEDIA_ROOT, 'translated_files', f'{os.path.basename(file_obj.name).replace(".", "_")}_translated.xlsx')
    workbook.save(translated_file_path)
    return translated_file_path



def translate_pptx(file_obj, target_language):
    try:
        prs = Presentation(file_obj)
    except Exception as e:
        return str(e)

    translate_and_preserve_formatting(prs, target_language)

    translated_file_path = os.path.join(settings.MEDIA_ROOT, 'translated_files', f'{os.path.basename(file_obj.name).replace(".", "_")}_translated.xlsx')
    prs.save(translated_file_path)
    return translated_file_path



def download_translated_file(request, file_id):
    translated_file = get_object_or_404(TranslatedFile, id=file_id)
    if not translated_file.translated_file:
        raise Http404("Translated file not available for download.")


    with open(translated_file.translated_file.path, 'rb') as f:
        filename = translated_file.original_file.name  # Use original filename for better clarity
        response = FileResponse(f, as_attachment=True)
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response


